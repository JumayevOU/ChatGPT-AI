"""extract_text_from_document() uchun qo'lda ishga tushiriladigan tekshiruv.
Ishga tushirish: python test_extract.py

Nega kerak: Excel/PPTX uchun o'quvchi bo'lmaganda kod faylni oddiy UTF-8
matn sifatida "dekodlab", GPT'ga ma'nosiz binar chiqindi yuborardi. Model
esa haqli ravishda "buzilgan binar matn yuborilgan, tahrirlay olmayman"
deb javob berardi — foydalanuvchi aynan shuni ko'rgan edi.
"""
import asyncio
import io

from services import extract_text_from_document, _looks_binary


def make_xlsx():
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws["A1"], ws["B1"] = "Ism", "Tug'ilgan yil"
    ws["A2"], ws["B2"] = "Ali", "31.12.99"
    ws["A3"], ws["B3"] = "Vali", "15.03.01"
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_xls():
    import xlwt
    wb = xlwt.Workbook()
    ws = wb.add_sheet("Sheet1")
    for r, row in enumerate([["Ism", "Tug'ilgan yil"], ["Ali", "31.12.99"]]):
        for c, v in enumerate(row):
            ws.write(r, c, v)
    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def make_pptx():
    from pptx import Presentation
    from pptx.util import Inches
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "Birinchi slayd matni"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


async def main():
    # 1) XLSX — o'qiladigan jadval ko'rinishi qaytishi kerak
    t = await extract_text_from_document(make_xlsx(), "maktab.xlsx")
    assert "31.12.99" in t, f"qiymat topilmadi: {t[:200]!r}"
    assert "Ism" in t and "Vali" in t, t[:200]
    assert not _looks_binary(t), "xlsx binar deb baholandi"
    print("[1] .xlsx o'qildi OK")

    # 2) XLS (eski BIFF) — aynan foydalanuvchi yuborgan format
    t = await extract_text_from_document(make_xls(), "maktab.xls")
    assert "31.12.99" in t, f"qiymat topilmadi: {t[:200]!r}"
    assert not t.startswith("[BINARY]"), "xls binar deb belgilandi"
    assert not _looks_binary(t), "xls chiqindi qaytardi"
    print("[2] .xls o'qildi OK")

    # 3) PPTX — slayd matnlari
    t = await extract_text_from_document(make_pptx(), "taqdimot.pptx")
    assert "Birinchi slayd matni" in t, t[:200]
    print("[3] .pptx o'qildi OK")

    # 4) Haqiqiy binar (noma'lum format) — [BINARY] deb belgilanishi kerak,
    #    chiqindi matn GPT'ga YUBORILMASLIGI shart
    blob = b"\x89PNG\r\n\x1a\n" + bytes(range(256)) * 30
    t = await extract_text_from_document(blob, "rasm.dat")
    assert t == "[BINARY]", f"binar belgilanmadi: {t[:120]!r}"
    print("[4] binar fayl [BINARY] deb belgilandi OK")

    # 5) Oddiy matn formatlari buzilmagan (regressiya)
    t = await extract_text_from_document("ism,ball\nAli,70\n".encode(), "a.csv")
    assert "Ali,70" in t, t
    t = await extract_text_from_document('{"a": 1}'.encode(), "a.json")
    assert '"a": 1' in t, t
    t = await extract_text_from_document("Salom dunyo".encode(), "a.txt")
    assert t == "Salom dunyo", repr(t)
    print("[5] csv/json/txt regressiyasi yo'q OK")

    # 6) _looks_binary chegaralari
    assert _looks_binary("") is True
    assert _looks_binary("oddiy o'zbekcha matn, hammasi joyida") is False
    assert _looks_binary("\x00\x01\x02\x03\x04" * 50) is True
    print("[6] _looks_binary OK")

    print("\nextract: barcha tekshiruvlar o'tdi (6/6).")


if __name__ == "__main__":
    asyncio.run(main())
