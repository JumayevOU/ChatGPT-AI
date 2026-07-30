import os
import shutil
import speech_recognition as sr
from pydub import AudioSegment
import edge_tts
import asyncio
import aiohttp
import re
import json
from contextlib import AsyncExitStack
from datetime import datetime, timezone, timedelta
from typing import List, Dict, Optional
from html import escape as html_escape

from ddgs import DDGS
from openai import NotFoundError, RateLimitError

if shutil.which("ffmpeg"):
    AudioSegment.converter = "ffmpeg"
else:
    AudioSegment.converter = "ffmpeg.exe"

from file_task_quota import FileTaskQuota
from sandbox import run_in_sandbox

try:
    from config import (
        GPT_MODEL, MODEL_FALLBACKS, CONTEXT_WINDOW, OPENAI_API_KEY,
        REQUEST_TIMEOUT, CONCISE_INSTRUCTION, STRICT_MATH_RULES,
        MESSAGE_COST_FILE_TASK,
        build_system_prompt, build_request_params, pick_reasoning_effort,
    )
except ImportError:
    import os
    OPENAI_API_KEY = os.getenv("OPENAI_API_KEY")
    CONCISE_INSTRUCTION = ""
    STRICT_MATH_RULES = ""
    GPT_MODEL = "gpt-4o-mini"
    MODEL_FALLBACKS = []
    REQUEST_TIMEOUT = 60.0
    CONTEXT_WINDOW = 12
    MESSAGE_COST_FILE_TASK = 250

    def build_system_prompt() -> str:
        return "You are a helpful assistant."

    def build_request_params(user_text: str = "", force_deep: bool = False, model: Optional[str] = None) -> Dict:
        return {"model": model or GPT_MODEL, "max_output_tokens": 1500, "reasoning": {"effort": "low"}}

    def pick_reasoning_effort(text: str, force_deep: bool = False) -> str:
        return "low"

from loader import openai_client, logger
from utils.history import update_chat_history

# ─────────────────────────────────────────────────────────────
# YORDAMCHI: TARIX FUNKSIYALARI
# ─────────────────────────────────────────────────────────────

async def clear_chat_history(chat_id: int):
    try:
        from utils.history import clear_history
        await clear_history(chat_id)
    except Exception as e:
        logger.error(f"Xotirani tozalashda xatolik: {e}")

async def safe_update_history(chat_id: int, content: str, role: str = "user"):
    if not content:
        return
    try:
        await update_chat_history(chat_id, content, role=role)
    except Exception as e:
        # Xabar HECH QAYERGA saqlanmagani keyingi javoblarda kontekst
        # yo'qolishiga bevosita olib keladi, shuning uchun warning.
        logger.warning(f"[Tarix yozish xatosi] chat={chat_id}, role={role}: {e}")

async def safe_get_chat_history(chat_id: int, limit: int = CONTEXT_WINDOW) -> List[Dict[str, str]]:
    try:
        from utils.history import get_chat_history
        hist = await get_chat_history(chat_id, limit=limit)
        return hist[-limit:] if isinstance(hist, list) else []
    except Exception as e:
        # DIQQAT: bu yerda xatolik "yutilib" bo'sh ro'yxat qaytarilsa,
        # foydalanuvchi uchun BUTUN suhbat konteksti indamay yo'qoladi.
        # Shuning uchun debug emas, warning darajasida logga yoziladi —
        # aks holda production'da bunday holatlarni payqash deyarli
        # imkonsiz bo'ladi.
        logger.warning(f"[Tarix o'qish xatosi] chat={chat_id}: {e} — bo'sh tarix qaytarilmoqda!")
        return []

# ─────────────────────────────────────────────────────────────
# RICH MESSAGE HELPERS
# ─────────────────────────────────────────────────────────────

_RICH_CODE_RE = re.compile(r"```([a-zA-Z0-9_+-]*)\n(.*?)\n```", re.S)
_RICH_MATH_BLOCK_RE = re.compile(r"(?<!\\)\$\$(.+?)\$\$", re.S)
_RICH_MATH_INLINE_RE = re.compile(r"(?<!\\)\$(?!\s)(.+?)(?<!\s)\$(?!\$)", re.S)
_RICH_DATE_PATTERNS = (
    re.compile(r"\b(\d{4})-(\d{2})-(\d{2})(?:[ T](\d{2}):(\d{2}))?\b"),
    re.compile(r"\b(\d{2})\.(\d{2})\.(\d{4})(?:[ T](\d{2}):(\d{2}))?\b"),
)


def _protect_code_blocks(text: str):
    placeholders = []

    def repl(match):
        token = f"@@CODE_BLOCK_{len(placeholders)}@@"
        placeholders.append((token, match.group(0)))
        return token

    return _RICH_CODE_RE.sub(repl, text), placeholders


def _restore_code_blocks(text: str, placeholders):
    for token, block in placeholders:
        text = text.replace(token, block)
    return text


def _looks_like_math(expr: str) -> bool:
    expr = expr.strip()
    if not expr:
        return False
    if len(expr) > 120:
        return False
    math_chars = set("=+-*/^_()[]{}<>∑∫√π∞≈≠≤≥·×÷")
    return any(ch in math_chars for ch in expr) or (any(ch.isdigit() for ch in expr) and any(ch.isalpha() for ch in expr))


def _replace_math_block(match):
    expr = match.group(1).strip()
    if not _looks_like_math(expr):
        return match.group(0)
    return f"<tg-math-block>{html_escape(expr)}</tg-math-block>"


def _replace_math_inline(match):
    expr = match.group(1).strip()
    if not _looks_like_math(expr):
        return match.group(0)
    return f"<tg-math>{html_escape(expr)}</tg-math>"


def _replace_dates(text: str) -> str:
    tz_uz = timezone(timedelta(hours=5))

    def fmt_datetime(dt: datetime, original: str, has_time: bool) -> str:
        unix = int(dt.timestamp())
        fmt = "wDT" if has_time else "d"
        return f'<tg-time unix="{unix}" format="{fmt}">{html_escape(original)}</tg-time>'

    def repl_iso(match):
        year, month, day, hour, minute = match.groups()
        has_time = hour is not None and minute is not None
        try:
            if has_time:
                dt = datetime(int(year), int(month), int(day), int(hour), int(minute), tzinfo=tz_uz)
            else:
                dt = datetime(int(year), int(month), int(day), 0, 0, tzinfo=tz_uz)
            return fmt_datetime(dt, match.group(0), has_time)
        except Exception:
            return match.group(0)

    def repl_eu(match):
        day, month, year, hour, minute = match.groups()
        has_time = hour is not None and minute is not None
        try:
            if has_time:
                dt = datetime(int(year), int(month), int(day), int(hour), int(minute), tzinfo=tz_uz)
            else:
                dt = datetime(int(year), int(month), int(day), 0, 0, tzinfo=tz_uz)
            return fmt_datetime(dt, match.group(0), has_time)
        except Exception:
            return match.group(0)

    text = _RICH_DATE_PATTERNS[0].sub(repl_iso, text)
    text = _RICH_DATE_PATTERNS[1].sub(repl_eu, text)
    return text


def build_rich_markdown(text: str) -> str:
    """Best-effort rich message markdown for Telegram Bot API 10.1."""
    if not text:
        return ""
    protected, placeholders = _protect_code_blocks(text)
    protected = _RICH_MATH_BLOCK_RE.sub(_replace_math_block, protected)
    protected = _RICH_MATH_INLINE_RE.sub(_replace_math_inline, protected)
    protected = _replace_dates(protected)
    protected = _restore_code_blocks(protected, placeholders)
    return protected

# ─────────────────────────────────────────────────────────────
# YOUTUBE SUMMARY
# ─────────────────────────────────────────────────────────────

async def get_youtube_summary(chat_id: int, video_id: str, user_prompt: str = ""):
    def _fetch_transcript():
        import youtube_transcript_api
        try:
            return youtube_transcript_api.YouTubeTranscriptApi.get_transcript(video_id, languages=['uz', 'ru', 'en'])
        except:
            try:
                transcript_list = youtube_transcript_api.YouTubeTranscriptApi.list_transcripts(video_id)
                for transcript in transcript_list:
                    return transcript.fetch()
            except Exception:
                return []
        return []

    try:
        transcript_data = await asyncio.to_thread(_fetch_transcript)
        
        if not transcript_data:
            yield "Kechirasiz, bu videoning ochiq subtitrlari yo'q ekan (yoki video yopiq/musiqiy). Boshqa video yuborib ko'ring."
            return
            
        full_text = " ".join([t.get('text', '') for t in transcript_data])
        
        if len(full_text) > 15000:
            full_text = full_text[:15000] + "\n...[Xarajatni tejash uchun videoning qolgan qismi qisqartirildi]."

        default_prompt = "Shu videoni qisqacha xulosa qilib ber va asosiy g'oyalarini ayt"
        final_prompt_text = user_prompt if user_prompt else default_prompt

        prompt = (
            f"Quyida YouTube videosining matni (subtitrlari) berilgan. "
            f"Foydalanuvchining so'rovi: '{final_prompt_text}'.\n\n"
            f"Video matni:\n{full_text}"
        )

        async for chunk in get_gpt_reply(chat_id, prompt):
            yield chunk

    except Exception as e:
        logger.error(f"YouTube Transcript xatosi: {e}")
        yield "Videoni tahlil qilishda kutilmagan xatolik yuz berdi."

# ─────────────────────────────────────────────────────────────
# 🔍 YAXSHILANGAN QIDIRUV BLOKI
# ─────────────────────────────────────────────────────────────

async def fetch_page_content(url: str, max_chars: int = 4000) -> str:
    """
    Berilgan URL dan sahifaning to'liq matnini yuklaydi va
    HTML teglarini olib tashlab, toza matn qaytaradi.
    """
    try:
        timeout = aiohttp.ClientTimeout(total=12)
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0 Safari/537.36"
            ),
            "Accept-Language": "uz,ru;q=0.9,en;q=0.8",
        }
        async with aiohttp.ClientSession(timeout=timeout) as session:
            async with session.get(url, headers=headers, ssl=False) as resp:
                if resp.status != 200:
                    return ""
                ct = resp.headers.get("Content-Type", "")
                if "text/html" not in ct and "text/plain" not in ct:
                    return ""
                html = await resp.text(errors="ignore")

                html = re.sub(r"<script[^>]*>.*?</script>", " ", html, flags=re.DOTALL | re.IGNORECASE)
                html = re.sub(r"<style[^>]*>.*?</style>",  " ", html, flags=re.DOTALL | re.IGNORECASE)
                html = re.sub(r"<nav[^>]*>.*?</nav>",       " ", html, flags=re.DOTALL | re.IGNORECASE)
                html = re.sub(r"<footer[^>]*>.*?</footer>", " ", html, flags=re.DOTALL | re.IGNORECASE)
                html = re.sub(r"<header[^>]*>.*?</header>", " ", html, flags=re.DOTALL | re.IGNORECASE)
                clean = re.sub(r"<[^>]+>", " ", html)
                clean = re.sub(r"&[a-zA-Z]{2,6};", " ", clean)
                clean = re.sub(r"&#\d+;", " ", clean)
                clean = re.sub(r"\s{2,}", " ", clean).strip()

                return clean[:max_chars]

    except asyncio.TimeoutError:
        logger.debug(f"fetch_page_content timeout: {url}")
    except Exception as e:
        logger.debug(f"fetch_page_content xatosi ({url}): {e}")
    return ""


async def multi_source_deep_search(
    primary_query: str,
    extra_queries: Optional[List[str]] = None,
    fetch_pages: int = 3,
) -> str:
    """
    Bir nechta so'rov bilan chuqur qidiruv:
      1. primary_query + extra_queries orqali qidiruv snippetlari
      2. Eng yuqori N ta URL dan to'liq sahifa matni yuklanadi
      3. Hammasi bitta kontekst sifatida qaytariladi
    """
    queries: List[str] = [primary_query]
    if extra_queries:
        queries += extra_queries[:2]   

    seen_urls: set = set()
    all_snippets: List[str] = []
    top_urls: List[tuple] = []      

    for q in queries:
        def _s(query=q):
            try:
                with DDGS() as ddgs:
                    return list(ddgs.text(query, max_results=6))
            except Exception as e:
                logger.error(f"DDGS error [{query}]: {e}")
                return []

        results = await asyncio.to_thread(_s)
        if not results:
            all_snippets.append(f"⚠️ «{q}» bo'yicha natija topilmadi.")
            continue

        block = f"📌 QIDIRUV: «{q}»\n{'─'*50}\n"
        for i, r in enumerate(results, 1):
            title = r.get("title", "")
            url   = r.get("href", "")
            body  = r.get("body", "")
            block += f"[{i}] {title}\n    🔗 {url}\n    {body}\n\n"

            if url and url not in seen_urls and i <= 2:
                seen_urls.add(url)
                top_urls.append((url, title))

        all_snippets.append(block)

    snippets_text = "\n\n".join(all_snippets)

    if top_urls and fetch_pages > 0:
        urls_to_fetch = top_urls[:fetch_pages]
        tasks = [fetch_page_content(url) for url, _ in urls_to_fetch]
        page_contents = await asyncio.gather(*tasks, return_exceptions=True)

        pages_block = "\n\n📄 SAHIFALARDAN TO'LIQ MA'LUMOT:\n" + "═" * 55 + "\n"
        any_page = False
        for (url, title), content in zip(urls_to_fetch, page_contents):
            if isinstance(content, str) and len(content) > 150:
                pages_block += f"\n🌐 {title}\n🔗 {url}\n\n{content[:3500]}\n{'─'*50}\n"
                any_page = True

        if any_page:
            snippets_text += pages_block

    return snippets_text if snippets_text.strip() else "Hech qanday ma'lumot topilmadi."


# ─────────────────────────────────────────────────────────────
# HUJJAT VA RASM
# ─────────────────────────────────────────────────────────────

def _extract_text_from_document_sync(file_bytes: bytes, filename: str) -> str:
    """
    Hujjatdan matnni ajratib olishning bloklovchi (CPU/IO) qismi.
    fitz (PDF) va python-docx kutubxonalari sinxron ishlaydi, shuning uchun
    bu funksiya faqat `asyncio.to_thread` orqali, alohida oqimda chaqiriladi —
    event loopni band qilmaslik uchun.
    """
    ext = filename.lower().split('.')[-1] if '.' in filename else ''
    text = ""

    if ext == 'pdf':
        # PDF fayllarni o'qish (eskisi kabi 10 sahifa limiti bilan)
        import fitz
        pdf_document = fitz.open(stream=file_bytes, filetype="pdf")
        pages_to_read = min(10, len(pdf_document))
        for page_num in range(pages_to_read):
            page = pdf_document.load_page(page_num)
            text += page.get_text()
        if len(pdf_document) > 10:
            text += "\n\n[TIZIM XABARI: Xarajat va xotirani tejash maqsadida hujjatning faqat dastlabki 10 sahifasi o'qildi.]"

    elif ext in ['doc', 'docx']:
        # Word fayllarini o'qish
        from io import BytesIO
        import docx
        doc = docx.Document(BytesIO(file_bytes))
        text = "\n".join([para.text for para in doc.paragraphs])

    elif ext in ['xlsx', 'xlsm', 'xls']:
        # Excel. DIQQAT: bularsiz fayl pastdagi `else` shoxiga tushib,
        # UTF-8 sifatida "dekodlanardi" va natijada GPT'ga ma'nosiz binar
        # chiqindi "Hujjat matni" nomi ostida yuborilardi. Model esa
        # haqli ravishda "buzilgan binar matn yuborilgan" deb javob berardi.
        text = _read_spreadsheet(file_bytes, ext)

    elif ext in ['pptx']:
        from io import BytesIO
        from pptx import Presentation
        prs = Presentation(BytesIO(file_bytes))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [
                sh.text_frame.text.strip()
                for sh in slide.shapes
                if sh.has_text_frame and sh.text_frame.text.strip()
            ]
            if texts:
                parts.append(f"--- Slayd {i} ---\n" + "\n".join(texts))
        text = "\n\n".join(parts)

    else:
        # Qolgan formatlar (txt, csv, json, kod fayllari, va h.k.) —
        # oddiy matn sifatida o'qishga harakat qilamiz.
        text = file_bytes.decode('utf-8', errors='ignore')
        if _looks_binary(text):
            # Binar fayl matn sifatida o'qilmaydi. Chiqindini GPT'ga
            # yubormaymiz — chaqiruvchi buni ko'rib, faylni sandbox
            # orqali ishlashini bildiradi.
            return "[BINARY]"

    # AI xotirasi to'lib qolmasligi uchun barcha fayllarga umumiy belgi cheklovi
    return text[:15000]


def _looks_binary(text: str) -> bool:
    """Dekodlangan matn aslida binar chiqindi ekanini aniqlaydi.

    Binar fayl `errors='ignore'` bilan dekodlanganda bo'sh emas, lekin
    ichida ko'p boshqaruv belgilari bo'ladi. Shu nisbatga qarab hukm
    qilamiz.
    """
    if not text:
        return True
    sample = text[:4000]
    weird = sum(1 for ch in sample if not (ch.isprintable() or ch.isspace()))
    return weird / len(sample) > 0.15


def _read_spreadsheet(file_bytes: bytes, ext: str) -> str:
    """Excel faylidan o'qiladigan ko'rinishdagi qisqa xulosa tayyorlaydi.

    Butun jadvalni emas, har varaqdan bosh qismini beradi — GPT'ga
    strukturani (ustun nomlari, qiymat ko'rinishi) tushunish uchun shu
    yetarli, tokenlar esa behuda sarflanmaydi.
    """
    from io import BytesIO

    MAX_ROWS, MAX_COLS = 25, 15
    out = []

    def fmt(rows, title):
        out.append(f"--- {title} ---")
        for row in rows:
            cells = ["" if c is None else str(c) for c in row[:MAX_COLS]]
            out.append(" | ".join(cells))

    if ext == 'xls':
        import xlrd
        book = xlrd.open_workbook(file_contents=file_bytes)
        for sheet in book.sheets():
            rows = [sheet.row_values(r) for r in range(min(MAX_ROWS, sheet.nrows))]
            fmt(rows, f"Varaq '{sheet.name}' ({sheet.nrows} qator x {sheet.ncols} ustun)")
    else:
        import openpyxl
        wb = openpyxl.load_workbook(BytesIO(file_bytes), data_only=True, read_only=True)
        for ws in wb.worksheets:
            rows = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= MAX_ROWS:
                    break
                rows.append(row)
            fmt(rows, f"Varaq '{ws.title}' ({ws.max_row} qator x {ws.max_column} ustun)")
        wb.close()

    return "\n".join(out)


async def extract_text_from_document(file_bytes: bytes, filename: str) -> str:
    """Hujjatlardan matnni ajratib olish (PDF, DOCX va barcha boshqa formatlar).

    Diqqat: PDF/DOCX parsing CPU/IO bo'yicha og'ir va sinxron kutubxonalar
    (fitz, python-docx) orqali amalga oshiriladi. Shuning uchun bu ishni
    to'g'ridan-to'g'ri bajarish event loopni bloklab, o'sha vaqtda BOSHQA
    barcha foydalanuvchilarning botga yuborgan xabarlarini "muzlatib" qo'yardi.
    `asyncio.to_thread` yordamida bu og'ir qism alohida oqimga o'tkazilgan.
    """
    try:
        return await asyncio.to_thread(_extract_text_from_document_sync, file_bytes, filename)
    except Exception as e:
        logger.error(f"Faylni o'qishda xatolik ({filename}): {e}")
        return "[XATOLIK] Ushbu fayl formatidan matnni ajratib olib bo'lmadi."


async def _open_response_stream(stack: AsyncExitStack, candidate_models: List[str], **kwargs):
    """
    `candidate_models` ro'yxatini ketma-ket sinaydi — asosiy model 404
    (mavjud emas) yoki 429 (limitga yetgan) qaytarsa, keyingi zaxira
    modelga o'tadi (config.py'dagi MODEL_FALLBACKS). Muvaffaqiyatli ochilgan
    oqimni (uni yopish `stack` orqali kafolatlanadi) va qaysi model
    ishlaganini qaytaradi.
    """
    last_err: Optional[Exception] = None
    for candidate in candidate_models:
        try:
            stream = await stack.enter_async_context(
                openai_client.responses.stream(model=candidate, timeout=REQUEST_TIMEOUT, **kwargs)
            )
            return stream, candidate
        except (NotFoundError, RateLimitError) as e:
            last_err = e
            logger.warning(
                f"Model '{candidate}' ishlamadi ({type(e).__name__}), keyingi zaxira modelga o'tilmoqda..."
            )
            continue
    raise last_err


async def get_vision_reply(chat_id: int, base64_image: str, user_message: str, *, model: str = GPT_MODEL):
    system_prompt = f"{build_system_prompt()}\n\n{CONCISE_INSTRUCTION}\n\n{STRICT_MATH_RULES}"

    messages: list = []

    # BONUS TUZATISH: avvalgi versiyada rasm tahlili har doim "nol"
    # kontekstdan boshlanardi — suhbat tarixi umuman qo'shilmasdi. Endi
    # matnli oqim bilan bir xil tamoyilda oldingi xabarlar ham kiritiladi,
    # shunda foydalanuvchi rasm haqida davom etuvchi savol bersa
    # ("bu yerdagi ikkinchi odam-chi?"), model buni oldingi rasm/suhbat
    # bilan bog'lay oladi.
    recent = await safe_get_chat_history(chat_id, limit=CONTEXT_WINDOW)
    for m in recent:
        if "role" in m and "content" in m:
            messages.append({"role": m["role"], "content": m["content"]})

    messages.append({
        "role": "user",
        "content": [
            {"type": "input_text", "text": user_message},
            {
                "type": "input_image",
                "image_url": f"data:image/jpeg;base64,{base64_image}",
                "detail": "auto",
            },
        ],
    })

    base_params = build_request_params(user_text=user_message, model=model)
    initial_model = base_params.pop("model")

    try:
        async with AsyncExitStack() as stack:
            stream, _resolved_model = await _open_response_stream(
                stack,
                [initial_model, *MODEL_FALLBACKS],
                input=messages,
                instructions=system_prompt,
                store=False,
                **base_params,
            )
            async for event in stream:
                if event.type == "response.output_text.delta":
                    yield event.delta
            await stream.get_final_response()
    except Exception as e:
        logger.error(f"Vision API xatosi: {e}")
        raise


# ─────────────────────────────────────────────────────────────
# ROLE DETECTION
# ─────────────────────────────────────────────────────────────

def detect_role_from_text(text: str) -> str:
    t = text.lower()
    tech    = ["kod", "error", "xato", "python", "javascript", "ai", "api", "server", "sql"]
    sales   = ["narx", "sotish", "savdo", "mijoz", "reklama", "marketing"]
    psycho  = ["ruhiy", "psixolog", "depress", "stress", "maslahat"]

    def _has_word(words: list) -> bool:
        # Chegara (\b) bilan qidiriladi — aks holda "api" so'zi "kaPItal"
        # ichida ham "topilib", bexosdan noto'g'ri uslub tanlanardi.
        return any(re.search(rf"\b{re.escape(w)}\b", t) for w in words)

    if _has_word(tech):   return "technical"
    if _has_word(sales):  return "commercial"
    if _has_word(psycho): return "supportive"
    return ""

def role_instruction(role: str) -> str:
    if role == "technical":   return "Javobni texnik uslubda, aniq kod misollari yoki buyruqlar bilan taqdim et."
    if role == "commercial":  return "Javobni tijoriy, qisqa va savdoga yo'naltirilgan tilda bering."
    if role == "supportive":  return "Javobni yumshoq, empatik va qo'llab-quvvatlovchi uslubda bering."
    return ""


# ─────────────────────────────────────────────────────────────
# 🌟 ASOSIY GPT JAVOBI — MUKAMMAL MULTI-ROUND SEARCH
# ─────────────────────────────────────────────────────────────

# Responses API tool sxemasi Chat Completions'dan farqli — "function" ichiga
# emas, to'g'ridan-to'g'ri tekis (flat) shaklda beriladi.
_TOOLS = [
    {
        "type": "function",
        "name": "internet_search",
        "description": (
            "Real vaqt ma'lumotlarini (ob-havo, valyuta kursi, yangiliklar, narxlar, "
            "mahsulotlar, sport natijalari va boshqa o'zgaruvchan faktlar) internetdan "
            "qidirish uchun. O'zbekiston kontekstida: valyuta uchun cbu.uz, ob-havo uchun "
            "meteo.uz yoki uzgidromet, yangiliklar uchun kun.uz / gazeta.uz ishlatilsin."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "primary_query": {
                    "type": "string",
                    "description": (
                        "Asosiy qidiruv so'rovi. Imkon qadar aniq va manbani "
                        "ko'rsatib yozing. Masalan: 'USD UZS kursi bugun cbu.uz 2025', "
                        "'Toshkent ob-havo bugun meteo.uz', 'site:kun.uz so'nggi yangiliklar'."
                    ),
                },
                "extra_queries": {
                    "type": "array",
                    "items": {"type": "string"},
                    "description": (
                        "Qo'shimcha 1-2 ta qidiruv so'rovi. Manbalarni solishtirish yoki "
                        "ma'lumotni kengaytirish uchun ishlatiladi. Masalan birinchi so'rov "
                        "o'zbekcha bo'lsa, ikkinchisi ruscha yoki inglizcha bo'lishi mumkin."
                    ),
                },
            },
            "required": ["primary_query"],
        },
        "strict": False,
    }
]

# Hujjat DIZAYNI bo'yicha qo'llanma. Bu tool tavsifining bir qismi bo'ladi
# (config.py'dagi umumiy SYSTEM promptga EMAS) — chunki u faqat fayl bilan
# ishlaydigan oqimlarda kerak, va tool ro'yxati barqaror bo'lgani uchun
# prompt caching orqali deyarli tekinga tushadi.
#
# Sabab: bu qo'llanmasiz model juda oddiy, "oq fon + qora matn" hujjatlar
# yasaydi (bitta katta xat boshiga tiqilgan "•" belgilari bilan). Quyidagi
# qoidalar sinovdan o'tgan — python-pptx + matplotlib bilan haqiqatan ham
# professional ko'rinishdagi natija beradi.
_DOC_DESIGN_GUIDE = (
    "═══ HUJJAT DIZAYNI — MAJBURIY ═══\n"
    "Foydalanuvchi hujjat/taqdimot so'raganda natija PROFESSIONAL "
    "ko'rinishda bo'lishi SHART. Oddiy oq fonli, faqat qora matnli natija "
    "QABUL QILINMAYDI.\n\n"
    "1) FORMATNI TO'G'RI TANLANG:\n"
    "   - 'prezentatsiya', 'taqdimot', 'slayd' → PPTX (python-pptx), "
    "16:9 (13.333 x 7.5 dyuym). PDF QILMANG, chunki PPTX tahrirlanadi.\n"
    "   - 'hisobot', 'hujjat', 'ariza', 'xat' → DOCX (python-docx).\n"
    "   - 'jadval', 'ro'yxat', 'hisob-kitob' → XLSX (openpyxl).\n"
    "   - PDF faqat foydalanuvchi ANIQ 'PDF' desa (reportlab).\n\n"

    "1a) PDF UCHUN `docgen` MODULINI ISHLATING — MAJBURIY. Ish papkasida "
    "tayyor `docgen` moduli bor (import docgen). U shrift va matn "
    "o'lchamlarini o'z ustiga oladi — ENG KO'P UCHRAYDIGAN IKKI XATONI "
    "(■ kvadratlar va matnning sahifadan chiqib ketishi) shu hal qiladi:\n"
    "     import docgen\n"
    "     REG, BOLD = docgen.register_fonts()   # oʻzbekcha ʻ ni qoʻllaydi\n"
    "     BLUE = docgen.hex_rgb('#005BAA')\n"
    "     # bir qatorli matn — sigʻmasa shrift avtomatik kichrayadi:\n"
    "     docgen.draw_fitted(c, '1992', x, y, w, font=BOLD, size=38,\n"
    "                        color=BLUE, align='center')\n"
    "     # koʻp qatorli matn — avtomatik oʻraladi, PASTKI y ni qaytaradi:\n"
    "     y2 = docgen.draw_para(c, uzun_matn, x, y, w, font=REG, size=12)\n"
    "     # keyingi blokni y2 dan PASTDA boshlang — ustma-ust tushmaydi!\n"
    "     # gorizontal ustunlar uchun masshtab (yorliqqa joy qoldiradi):\n"
    "     k = docgen.bar_scale([9, 6, 15], maxw=4*inch, reserve=0.5*inch)\n"
    "   Boshqa funksiyalar: docgen.fit(), docgen.wrap(), "
    "docgen.para_height() (blokni joylashdan oldin balandligini bilish).\n"
    "   Helvetica ISHLATMANG — unda 'ʻ' yo'q va ■ bo'lib chiqadi.\n\n"

    "1b) PPTX/DOCX uchun shrift: 'Arial', 'Calibri', 'Segoe UI' yoki "
    "'Tahoma'. 'Verdana' ISHLATMANG — unda 'ʻ' yo'q. Matn qutisiga "
    "tf.word_wrap = True qo'ying, qutini matnga yetadigan qilib "
    "kengaytiring va uzun matnda shriftni kichraytiring.\n\n"

    "1c) BLOKLAR USTMA-UST TUSHMASIN: har bir blokdan keyin keyingisining "
    "y koordinatasini HISOBLANG (draw_para qaytargan y yoki "
    "para_height()). Koordinatalarni taxminan qo'ymang — matn quyidagi "
    "karta yoki diagramma ustiga chiqib ketadi.\n\n"

    "1d) IERARXIYA — stat kartada RAQAM katta, IZOH kichik. Teskarisi XATO:\n"
    "   ✓ '1992' 36-40pt qalin asosiy rangda, ostida 'Tashkil topgan' "
    "12-13pt kulrang\n"
    "   ✗ '1992' kichik, 'Tashkil topgan' katta — bu XATO va sig'maydi\n"
    "   Izoh 1-2 so'zdan oshmasin.\n\n"

    "1e) VERTIKAL KOMPOZITSIYA: kontent slaydning faqat o'rta tasmasiga "
    "tiqilib, yuqori va past qismi bo'sh qolmasin. Sarlavha yuqoridan "
    "~0.8 dyuym, asosiy blok balandlikning 55-70%ini egallasin, pastda "
    "kolontitul. Bo'sh joy ataylab va muvozanatli bo'lsin.\n\n"

    "2) RANG SXEMASI: mavzuga mos 1 ta asosiy + 1 ta urg'u rangi tanlang "
    "va butun hujjatda izchil ishlating (masalan aviakompaniya uchun ko'k "
    "#005BAA + to'q sariq #E88B00; moliya uchun to'q yashil; tibbiyot "
    "uchun ko'k-yashil). Neytral: #0A1F33 matn, #5A6B7B ikkilamchi matn, "
    "#F2F6FA och fon.\n\n"
    "3) SLAYD TURLARINI ARALASHTIRING (hammasi bir xil bo'lmasin):\n"
    "   - Sarlavha slaydi: BUTUN slayd asosiy rangga bo'yalgan, katta oq "
    "sarlavha (44-54pt) + kichik ochroq izoh.\n"
    "   - Infografika: 3-4 ta 'stat karta' — och fonli to'rtburchak, "
    "ustida rangli chiziq, ichida KATTA raqam (36-44pt, asosiy rangda) "
    "va ostida kichik izoh (12-14pt).\n"
    "   - Diagramma slaydi: chap tomonda och fonli ustun (matn), o'ngda "
    "diagramma rasmi.\n"
    "   - Vaqt chizig'i: to'q fon, gorizontal chiziq, ustida rangli "
    "doiralar + yil + voqea.\n"
    "   - Xulosa: to'q yoki rangli fon, qisqa kuchli gap.\n\n"
    "4) DIAGRAMMA (kamida 1 ta bo'lsin, ma'lumot bo'lsa): matplotlib "
    "bilan chizing → PNG qilib saqlang (dpi=200, transparent=True, "
    "bbox_inches='tight') → slaydga add_picture bilan qo'ying. Donut "
    "diagramma uchun wedgeprops=dict(width=0.42, edgecolor='white', "
    "linewidth=2). Diagrammani ish papkasiga saqlang (output/ ga EMAS — "
    "u yerga faqat yakuniy hujjat tushsin).\n\n"
    "5) TIPOGRAFIYA: sarlavha 40-54pt qalin, slayd sarlavhasi 30-34pt "
    "qalin, matn 13-16pt. Bir slaydda 40 SO'ZDAN OSHMASIN.\n\n"
    "6) RO'YXATLAR: har bir band ALOHIDA xat boshi (paragraph) bo'lsin. "
    "Bandlarni ';' bilan bitta xat boshiga tiqib qo'ymang — bu eng "
    "ko'p uchraydigan xato.\n\n"
    "7) TEXNIK: shakllarda s.line.fill.background() (chegara chizig'ini "
    "olib tashlaydi) va s.shadow.inherit = False (keraksiz soyani "
    "o'chiradi) ishlating. Matn va fon o'rtasida kuchli kontrast bo'lsin — "
    "och fonda och kulrang matn o'qilmaydi.\n\n"
    "8) BELGILAR: ✈ kabi emoji shakl ichida mayda va noaniq chiqadi — "
    "ular o'rniga rangli doira/to'rtburchak va matn yorliqlaridan "
    "foydalaning.\n\n"
    "9) RASM: sandbox internetga chiqmaydi, shuning uchun tayyor "
    "fotosurat YUKLAB OLIB BO'LMAYDI. Buning o'rniga rangli bloklar, "
    "diagrammalar, katta tipografika va geometrik shakllardan "
    "foydalaning. Foydalanuvchi 'rasm qo'sh' desa — nega qo'sha "
    "olmasligingizni qisqa tushuntiring va o'zi rasm yuborsa "
    "joylashtira olishingizni ayting."
)

# Fayl yaratish/tahrirlash tool'i — `_TOOLS`dan ALOHIDA saqlanadi, chunki u
# faqat natijani yetkazib bera oladigan oqimlarda (output_files ro'yxati
# berilganda) biriktiriladi. Guest rejimda, masalan, biriktirilmaydi —
# u yerda hujjat yuborib bo'lmaydi, shuning uchun tool'ni ko'rsatish faqat
# behuda kod bajarilishiga olib kelardi.
_FILE_TASK_TOOL = {
    "type": "function",
    "name": "run_python_sandbox",
    "description": (
        "Fayl YARATISH yoki foydalanuvchi yuborgan faylni TAHRIRLASH kerak "
        "bo'lganda Python kodi yozib, shu tool orqali bajaring. Masalan: "
        "matndan PDF/Word/Excel yaratish, Excel katakchalarini o'zgartirish, "
        "CSV'ni saralash yoki filtrlash, diagramma (grafik) chizish, "
        "hujjatlarni birlashtirish, ZIP arxiv yasash, formatdan formatga "
        "o'girish.\n\n"
        "QACHON ISHLATMASLIK KERAK: agar foydalanuvchi shunchaki fayl "
        "haqida savol bersa yoki mazmunini so'rasa — bu tool KERAK EMAS, "
        "oddiy matnli javob bering.\n\n"
        "MUHIT:\n"
        "- Kod joriy papkada (cwd) ishlaydi. Foydalanuvchi fayl yuborgan "
        "bo'lsa, u shu papkada `input.<kengaytma>` nomi bilan turadi "
        "(masalan input.xlsx, input.csv, input.pdf).\n"
        "- Natija fayl(lar)ini ALBATTA `output/` papkasiga yozing — faqat "
        "o'sha papkadagi fayllar foydalanuvchiga yuboriladi.\n"
        "- Fayl nomini mazmunli qo'ying (masalan output/hisobot.pdf).\n"
        "- Mavjud kutubxonalar: pandas, openpyxl (.xlsx o'qish/yozish), "
        "xlrd (.xls O'QISH), xlwt + xlutils (.xls YOZISH), python-docx, "
        "python-pptx, pypdf, reportlab, matplotlib, PyMuPDF (fitz), "
        "beautifulsoup4, lxml va standart kutubxona (json, csv, zipfile, "
        "sqlite3, xml, re). Internetga chiqish YO'Q — hech narsa yuklab "
        "olmang va pip install qilmang. Ro'yxatda YO'Q kutubxonani "
        "import qilmang.\n"
        "- ESKI .xls FAYLNI TAHRIRLASH: openpyxl uni OCHA OLMAYDI. To'g'ri "
        "yo'l — xlrd bilan o'qib, xlutils.copy bilan nusxa olib, xlwt "
        "bilan yozish:\n"
        "     import xlrd\n"
        "     from xlutils.copy import copy as xl_copy\n"
        "     rb = xlrd.open_workbook('input.xls', formatting_info=True)\n"
        "     wb = xl_copy(rb)\n"
        "     ws = wb.get_sheet(0)\n"
        "     ws.write(qator, ustun, yangi_qiymat)\n"
        "     wb.save('output/natija.xls')\n"
        "  (formatting_info=True ishlamasa, uni olib tashlab qayta uring.)\n"
        "- Kirill/lotin matn uchun faylni har doim encoding='utf-8' bilan "
        "yozing. PDF'da o'zbekcha harflar kerak bo'lsa, reportlab'ning "
        "o'rnatilgan Helvetica shrifti lotin harflarini qo'llab-quvvatlaydi.\n"
        "- Kod tugagach print() bilan qisqacha nima qilganingizni yozing.\n\n"
        "Fayl tuzilishini bilmasangiz, avval uni tekshiradigan qisqa kod "
        "yuboring (masalan sheet nomlari va birinchi qatorlarni print "
        "qiling), natijani ko'ring, keyin ikkinchi chaqiriqda haqiqiy "
        "o'zgartirishni bajaring — bu tool bir xabar davomida bir necha "
        "marta chaqirilishi mumkin. Kod xato bersa, xato matnini o'qib "
        "tuzating va qayta chaqiring.\n\n"
        + _DOC_DESIGN_GUIDE
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "code": {
                "type": "string",
                "description": "Bajariladigan to'liq Python kodi.",
            },
        },
        "required": ["code"],
    },
    "strict": False,
}

_SYNTHESIS_SYSTEM = """QAT'IY BUYRUQ — ANIQ VA CHUQUR JAVOB YOZ:

0. TIL — ENG MUHIM QOIDA: Yakuniy javobni albatta foydalanuvchining ASL savoli
   (yuqoridagi "user" xabari) qanday tilda yozilgan bo'lsa, AYNAN o'sha tilda yoz —
   bu qidiruv natijalari/manbalar qaysi tilda bo'lishidan (ular ko'pincha o'zbek
   tilida chiqadi, chunki qidiruv so'zlari cbu.uz/meteo.uz kabi manbalar uchun
   o'zbekcha tuziladi) qat'iy nazar. Masalan foydalanuvchi ruscha yozgan bo'lsa,
   manbalar o'zbekcha bo'lsa ham, javob albatta RUSCHA bo'lishi SHART. Bu qoidani
   hech qachon buzma — pastdagi format qoidalari ham shu tilda qo'llanadi.

1. MANBALARNI TAHLIL QIL: Berilgan barcha manba matnlarini o'qib, ularni SOLISHTIR.
   Bir manba boshqasiga zid bo'lsa — bu ziddiyatni foydalanuvchiga ayt.

2. MANTIQ: Harorat, kurs yoki raqamlarni HISOBLASHDA xato qilma.
   - Agar 24°C va yomg'ir kutilsa → qalin kiyim TAVSIYA ETMA.
   - Dollar kursi so'ralsa → faqat UZS qiymatini yoz, taxmin qilma.

3. FORMAT — qisqa va lo'nda TAQIQLANADI:
   - Sarlavhalar (bold) ishlat.
   - Raqamli yoki belgili ro'yxatlar tuz.
   - Tegishli emojilar qo'sh (☀️ 🌧️ 💰 📈 📰).
   - Kamida 3-5 xat boshi yoz.

4. MANBA: Javob oxirida foydalanilgan manbalar QU'YIDAGI FORMATDA ko'rsat:
   - [Sayt nomi](URL)
   Misol: - [Weather.com](https://weather.com/...) yoki - [CoinMarketCap](https://coinmarketcap.com/...)
   Hech qachon raw URL yozma. Doim [Nom](url) formatidan foydalan.
   Agar rasmiy sayt topilmasa — buni ochiq ayt.

5. SANA/VAQT: Agar ma'lumot eskirgan bo'lsa yoki aniq sana topilmasa — buni ham ayt."""


async def _run_file_task(
    code: str,
    *,
    quota: Optional[FileTaskQuota],
    input_file_bytes: Optional[bytes],
    input_filename: Optional[str],
    output_files: Optional[list],
    round_num: int,
    rounds_left: int = 99,
) -> str:
    """run_python_sandbox tool chaqiruvini bajaradi.

    Kvotani (bir marta) yechadi, kodni sandbox'da ishga tushiradi va
    natijani GPT tushunadigan matn sifatida qaytaradi. Yaratilgan
    fayllarni chaqiruvchining `output_files` ro'yxatiga qo'shadi.
    """
    if quota is not None:
        allowed = await quota.ensure_charged()
        if not allowed:
            return (
                "XATOLIK: foydalanuvchining bugungi krediti tugagan, kod "
                "bajarilmadi. Boshqa urinmang — foydalanuvchiga kredit "
                "tugaganini muloyim tushuntiring."
            )

    logger.info(f"[FileTask] round={round_num}, kod uzunligi={len(code)}")
    result = await run_in_sandbox(code, input_file_bytes, input_filename)

    # Model bosqichlar tugayotganini BILISHI kerak — aks holda oxirgi
    # urinishni ham tekshiruvga sarflab, keyin "menda bunday vosita yo'q"
    # deb noto'g'ri javob yozadi.
    warn = ""
    if rounds_left <= 1:
        warn = (
            "\n\nDIQQAT: bu SO'NGGI urinish edi, tool boshqa chaqirilmaydi. "
            "Agar fayl yaratilmagan bo'lsa — foydalanuvchiga muammoni "
            "aniq va qisqa tushuntiring ('vosita yo'q' DEMANG, chunki "
            "vosita bor edi)."
        )

    if not result.success:
        logger.info(f"[FileTask] round={round_num} XATO: {result.traceback[:200]}")
        return (
            f"XATO — kod bajarilmadi:\n{result.traceback[:3000]}\n\n"
            "Kodni tuzatib qayta chaqiring." + warn
        )

    if not result.output_files:
        return (
            f"Kod xatosiz bajarildi, LEKIN `output/` papkasida hech qanday "
            f"fayl yo'q.\nSTDOUT:\n{result.stdout[:1500]}\n\n"
            "Agar bu tekshiruv (inspeksiya) qadami bo'lsa — davom eting va "
            "endi haqiqiy faylni yarating. Agar fayl yaratmoqchi bo'lgan "
            "bo'lsangiz — uni `output/` papkasiga yozganingizga ishonch "
            "hosil qiling." + warn
        )

    if quota is not None:
        quota.mark_success()
    if output_files is not None:
        output_files.extend(result.output_files)

    names = ", ".join(f"{n} ({len(b)} bayt)" for n, b in result.output_files)
    logger.info(f"[FileTask] round={round_num} muvaffaqiyat: {names}")
    return (
        f"BAJARILDI. Fayllar yaratildi va foydalanuvchiga avtomatik "
        f"yuboriladi: {names}\nSTDOUT:\n{result.stdout[:1500]}\n\n"
        "Endi foydalanuvchiga nima qilganingizni QISQA (1-2 gap) tushuntiring. "
        "Faylni qanday yuklab olishni tushuntirmang — u allaqachon biriktirilgan."
    )


async def get_openai_reply(
    chat_id: int,
    message_text: str,
    *,
    model: str = GPT_MODEL,
    user_id: Optional[int] = None,
    input_file_bytes: Optional[bytes] = None,
    input_filename: Optional[str] = None,
    output_files: Optional[list] = None,
):
    system_prompt = f"{build_system_prompt()}\n\n{CONCISE_INSTRUCTION}\n\n{STRICT_MATH_RULES}"

    messages: list = []

    try:
        now_utc = datetime.now(timezone.utc)
        now_tashkent = now_utc.astimezone(timezone(timedelta(hours=5)))
        time_msg = (
            f"[TIZIM MA'LUMOTI]\n"
            f"Hozirgi sana: {now_tashkent.strftime('%Y-%m-%d')}, "
            f"Vaqt: {now_tashkent.strftime('%H:%M')} (O'zbekiston, UTC+5).\n"
            f"Foydalanuvchi O'zbekistonda. 'Dollar' = USD/UZS kursi (Markaziy bank). "
            f"'Ob-havo' = O'zbekiston hududi (uzgidromet / meteo.uz).\n"
            f"Real vaqt ma'lumotlari uchun DOIM 'internet_search' asbobini ishlat — "
            f"o'z bilimingdan javob to'qima!"
        )
        messages.append({"role": "developer", "content": time_msg})
    except Exception:
        pass

    recent = await safe_get_chat_history(chat_id, limit=CONTEXT_WINDOW)
    for m in recent:
        if "role" in m and "content" in m:
            messages.append({"role": m["role"], "content": m["content"]})

    role = detect_role_from_text(message_text)
    r_instr = role_instruction(role)
    if r_instr:
        messages.append({"role": "developer", "content": f"ROLE_INSTRUCTION: {r_instr}"})

    messages.append({"role": "user", "content": message_text})

    # Reasoning effort xabar murakkabligiga qarab tanlanadi (config.py:
    # pick_reasoning_effort) — soddasiga tez/arzon, murakkabiga chuqurroq.
    base_params = build_request_params(user_text=message_text, model=model)
    initial_model = base_params.pop("model")

    # Qidiruv va fayl vazifasi uchun ALOHIDA byudjet. Avval ikkalasi bitta
    # 3 bosqichli hisobni bo'lishardi — natijada "qidirib, keyin hujjat
    # yasa" so'rovida qidiruv bosqichlarni yeb qo'yib, faylni yaratishga
    # urinish yetmay qolardi va foydalanuvchi javob olsa-da, FAYLSIZ qolardi.
    MAX_SEARCH_ROUNDS = 3
    # Fayl uchun 4: model odatda birinchi bosqichni fayl tuzilishini
    # TEKSHIRISHGA sarflaydi (bu to'g'ri xatti-harakat), keyin yozadi va
    # xato bo'lsa 1-2 marta tuzatadi. 3 ta bo'lganda tekshiruv + bitta
    # xato butun byudjetni yeb qo'yardi.
    MAX_FILE_ROUNDS = 4
    MAX_TOTAL_ROUNDS = 8   # cheksiz siklga qarshi umumiy xavfsizlik chegarasi

    search_rounds = 0
    file_rounds = 0
    total_rounds = 0
    search_performed = False
    synthesis_injected = False
    resolved_model: Optional[str] = None

    # Fayl tool'i faqat natijani yetkaza oladigan oqimlarda biriktiriladi
    # (chaqiruvchi `output_files` ro'yxatini bergan bo'lsa). Guest rejimda
    # bu None bo'ladi — u yerda hujjat yuborib bo'lmaydi.
    file_task_enabled = output_files is not None
    file_quota: Optional[FileTaskQuota] = (
        FileTaskQuota(user_id, MESSAGE_COST_FILE_TASK)
        if (file_task_enabled and user_id is not None)
        else None
    )
    file_task_started = False

    while True:
        # MUHIM: qidiruv 1-2 bosqichda tugasa ham (model ko'proq tool
        # so'ramasa), keyingi chaqiruvda hali ham `tools` biriktirilgan
        # bo'lishi mumkin (model xohlasa yana qidirishi mumkin) — shuning
        # uchun _SYNTHESIS_SYSTEM'ni majburiy yakuniy bosqichga emas,
        # birinchi qidiruv natijasi qaytgan zahoti (pastda) bir marta
        # qo'shamiz — u qachon javob bersa ham ishlaydi.
        active_tools = []
        if search_rounds < MAX_SEARCH_ROUNDS:
            active_tools.extend(_TOOLS)
        if file_task_enabled and file_rounds < MAX_FILE_ROUNDS:
            active_tools.append(_FILE_TASK_TOOL)

        call_kwargs = dict(base_params)
        call_kwargs.update(input=messages, instructions=system_prompt, store=False)
        if active_tools and total_rounds < MAX_TOTAL_ROUNDS:
            call_kwargs.update(tools=active_tools, tool_choice="auto")

        candidate_models = [resolved_model] if resolved_model else [initial_model, *MODEL_FALLBACKS]

        try:
            async with AsyncExitStack() as stack:
                stream, resolved_model = await _open_response_stream(stack, candidate_models, **call_kwargs)

                got_function_call = False
                pending_calls = []

                async for event in stream:
                    et = event.type
                    if et == "response.output_text.delta":
                        yield event.delta
                    elif et == "response.output_item.added" and getattr(event.item, "type", None) == "function_call":
                        got_function_call = True
                        # Kontent emas — faqat "band" animatsiyasiga signal:
                        # qidiruv yoki fayl vazifasi (sekundlab davom etadi)
                        # boshlanmoqda.
                        if getattr(event.item, "name", None) == "run_python_sandbox":
                            if not file_task_started:
                                yield "[STATUS]file_task"
                                file_task_started = True
                        elif not search_performed:
                            yield "[STATUS]search"
                            search_performed = True
                    elif et == "response.output_item.done" and getattr(event.item, "type", None) == "function_call":
                        pending_calls.append(event.item)

                final_response = await stream.get_final_response()
        except Exception as e:
            logger.error(f"GPT javob xatosi: {e}")
            raise

        if not got_function_call:
            if file_quota is not None:
                await file_quota.refund_if_unused()
            if final_response.status == "incomplete":
                logger.warning(f"GPT javobi incomplete tugadi: {final_response.incomplete_details}")
            return

        file_task_ran = False
        search_ran = False

        for call_item in pending_calls:
            try:
                args = json.loads(call_item.arguments or "{}")
            except Exception:
                args = {}

            if call_item.name == "run_python_sandbox":
                file_task_ran = True
                tool_output = await _run_file_task(
                    args.get("code", ""),
                    quota=file_quota,
                    input_file_bytes=input_file_bytes,
                    input_filename=input_filename,
                    output_files=output_files,
                    round_num=file_rounds + 1,
                    rounds_left=MAX_FILE_ROUNDS - file_rounds,
                )
            else:
                search_ran = True
                primary_query = args.get("primary_query", "")
                extra_queries = args.get("extra_queries", [])

                if primary_query:
                    logger.info(
                        f"[SEARCH] primary='{primary_query}' extra={extra_queries} round={search_rounds + 1}"
                    )
                    tool_output = await multi_source_deep_search(
                        primary_query=primary_query,
                        extra_queries=extra_queries if extra_queries else None,
                        fetch_pages=3,
                    )
                else:
                    tool_output = "Qidiruv so'rovi bo'sh bo'lgani uchun bajarilmadi."

            messages.append({
                "type": "function_call",
                "call_id": call_item.call_id,
                "name": call_item.name,
                "arguments": call_item.arguments,
            })
            messages.append({
                "type": "function_call_output",
                "call_id": call_item.call_id,
                "output": tool_output,
            })

        if search_ran:
            search_rounds += 1
        if file_task_ran:
            file_rounds += 1
        total_rounds += 1

        if file_task_ran:
            # Tool'dan OLDIN yozilgan oraliq matn ("Hozir tayyorlab
            # beraman...") yakuniy javobga yopishib qolmasligi uchun
            # ekranni tozalaymiz — yakuniy javob toza boshlanadi.
            yield "[CLEAR_TEXT]"

            # `_SYNTHESIS_SYSTEM` faqat INTERNET QIDIRUVI natijalarini
            # formatlash uchun (manbalar ro'yxati, emoji, kamida 3-5 xat
            # boshi). Fayl vazifasida u mutlaqo noo'rin — javob qisqa
            # bo'lishi kerak, chunki asosiy natija biriktirilgan faylning
            # o'zi. Shuning uchun uni qo'shmasdan keyingi bosqichga o'tamiz.
            continue

        if not synthesis_injected:
            # Birinchi qidiruv natijasi endigina qo'shildi — shu joydan
            # boshlab, model qachon (shu bosqichda yoki keyingisida)
            # yakuniy javob bersa ham, u manba-formatlash qoidasiga rioya
            # qiladi.
            messages.append({"role": "developer", "content": _SYNTHESIS_SYSTEM})
            yield "[CLEAR_TEXT]"
            synthesis_injected = True


async def get_gpt_reply(
    chat_id: int,
    user_message: str,
    *,
    user_id: Optional[int] = None,
    input_file_bytes: Optional[bytes] = None,
    input_filename: Optional[str] = None,
    output_files: Optional[list] = None,
):
    async for chunk in get_openai_reply(
        chat_id,
        user_message,
        user_id=user_id,
        input_file_bytes=input_file_bytes,
        input_filename=input_filename,
        output_files=output_files,
    ):
        yield chunk

# ─────────────────────────────────────────────────────────────
# STT, TTS
# ─────────────────────────────────────────────────────────────

# Bir xil STT servisi (Google recognize_google) bilan bir nechta til
# gipotezasini ketma-ket sinaymiz — bepul recognize_google bitta chaqiruvda
# tilni o'zi avtomatik aniqlay olmaydi, shuning uchun ovoz o'zbekcha bo'lmasa
# oldin "uz-UZ" xato/bo'sh natija berardi. Model/servis o'zgarmaydi, faqat
# shu servisga qaysi til bilan murojaat qilishni o'zimiz tanlaymiz.
_STT_LANGUAGE_CANDIDATES = ("uz-UZ", "ru-RU", "en-US")


def _speech_to_text_sync(file_path: str, wav_path: str) -> str:
    """
    Ovozni matnga aylantirishning bloklovchi qismi:
    - AudioSegment (pydub/ffmpeg) orqali formatni konvertatsiya qilish,
    - r.recognize_google() orqali Google Speech API'ga SINXRON tarmoq so'rovi.
    Bu ikkalasi ham event loopni sekundlab bloklashi mumkin, shuning uchun
    faqat `asyncio.to_thread` orqali, alohida oqimda chaqiriladi.
    """
    r = sr.Recognizer()
    audio = AudioSegment.from_file(file_path)
    audio.export(wav_path, format="wav")
    with sr.AudioFile(wav_path) as source:
        audio_data = r.record(source)

    for lang in _STT_LANGUAGE_CANDIDATES:
        try:
            text = r.recognize_google(audio_data, language=lang)
            if text:
                return text
        except sr.UnknownValueError:
            continue
    return ""


async def speech_to_text(file_path: str) -> str:
    """Ovozli xabarni matnga o'giradi.

    Diqqat: avvalgi versiyada butun konvertatsiya + Google'ga tarmoq so'rovi
    to'g'ridan-to'g'ri `async def` ichida, `await`siz bajarilardi — bu esa
    o'sha soniyalarda BUTUN botni (barcha foydalanuvchilar uchun) muzlatib
    qo'yardi. Endi bu og'ir ish `asyncio.to_thread` orqali alohida oqimda
    ishlaydi, event loop esa shu vaqtda boshqa foydalanuvchilarga xizmat
    qilishda davom etaveradi.
    """
    wav_path = file_path + ".wav"
    try:
        return await asyncio.to_thread(_speech_to_text_sync, file_path, wav_path)
    except Exception:
        return ""
    finally:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
            if os.path.exists(wav_path):
                os.remove(wav_path)
        except Exception:
            pass


async def text_to_speech(text: str, filename: str) -> str:
    text = text.replace("'", "'").replace("`", "'")
    clean_text_for_speech = re.sub(r'<[^>]+>', '', text)
    clean_text_for_speech = clean_text_for_speech.replace("$$", "").replace("$", "")
    VOICE = "uz-UZ-MadinaNeural"
    try:
        communicate = edge_tts.Communicate(clean_text_for_speech, VOICE, rate="-10%")
        await communicate.save(filename)
        return filename
    except Exception as e:
        logger.error(f"TTS xatosi: {e}")
        return None
